## one version for 8020

import pptx
import datetime
from openai import OpenAI
import base64
import requests
import json
import os
import random
from dotenv import load_dotenv
from io import BytesIO
from utils import WORK_FOLDER

load_dotenv()

client = OpenAI()


def gpt4_orchestrator(prompt, model="gpt-4-1106-preview"):
    response = client.chat.completions.create(
        model=model,
        response_format={"type": "json_object"},
        messages=[{"role": "user", "content": prompt}],
    )
    result = response.choices[0].message.content
    return result


def slide_skeleton(num_slides):
    """based on no. of slides create structure"""
    num_slides = max(3, min(int(num_slides), 5))
    slides = [
        {
            "slide_number": 1,
            "description": "title-slide",
            "title": "write here title very very short & succint",
            "subtitle": "write here subtitle of the entire presentation",
        }
    ]
    slides.append(
        {
            "slide_number": 2,
            "description": "intro-slide",
            "title": "write here intro title of this presentation followed by line break and subtitle",
            "undertitle": "write the main point of the entire presentation",
        }
    )
    for i in range(3, num_slides + 1):
        slides.append(
            {
                "slide_number": i,
                "description": "body-slide",
                "title": "write title of this slide very succintly followed by line break and subtitle",
                "object": "write here comprehensive answer to user query not just slide introduction but go deep using additional information be comprehensive and in one string max 500 tokens no line breaks",
            }
        )

    return {"slides": slides}


def create_prs(user_query, num_slides, context):
    presentation_slides = []
    structure = slide_skeleton(num_slides)
    for slide_info in structure["slides"]:
        slide_prompt = f"""
        Based on the user query: '{user_query}'\n
        and additional information: '{context}',\n
        create content and populate it into this JSON structure for the current slide: {slide_info}\n
        remember to write in the language of user query.
        """
        print(slide_prompt)
        try:
            slide_content_response = gpt4_orchestrator(slide_prompt)
            print(slide_content_response)
            slide_content = json.loads(slide_content_response)
            presentation_slides.append(slide_content)
        except Exception as e:
            print(f"error:{e}")
            continue

    return presentation_slides


def dalle_3(user_query):
    image_response = client.images.generate(
        model="dall-e-3",
        prompt=f"create beautiful photorealistic image taking into account context of user query:"
        + user_query,
        n=1,
        size="1024x1024",
        response_format="b64_json",
    )

    base_64 = base64.b64decode(image_response.data[0].b64_json)

    return base_64


def get_random_image(image_folder):
    images = [
        img
        for img in os.listdir(image_folder)
        if img.lower().endswith((".png", ".jpg", ".jpeg", ".gif"))
    ]
    if images:
        return os.path.join(image_folder, random.choice(images))
    return None


def generate_prs(user_id, user_query, num_slides, context=None):
    image_folder = "deck_generator/fallback_images"
    prs = pptx.Presentation("deck_generator/master.pptx")
    responses = create_prs(user_query, num_slides, context)
    presentation_title = (
        responses[0].get("title", "untitled") if responses else "untitled"
    )
    try:
        image_b64 = dalle_3(user_query)
        image_stream = BytesIO(image_b64)
    except Exception as e:
        random_image_path = get_random_image(image_folder)
        image_stream = random_image_path

    for response in responses:
        slide_data = response

        has_title = "title" in slide_data
        has_subtitle = "subtitle" in slide_data
        has_undertitle = "undertitle" in slide_data
        has_object = "object" in slide_data

        if has_title and has_subtitle:
            slide_layout = prs.slide_layouts[0]
        elif has_title and has_undertitle:
            slide_layout = prs.slide_layouts[1]
        elif has_title and has_object:
            slide_layout = prs.slide_layouts[2]

        slide = prs.slides.add_slide(slide_layout)

        if slide_layout == prs.slide_layouts[0]:
            slide.placeholders[1].text = slide_data.get("title", "")
            slide.placeholders[2].text = slide_data.get("subtitle", "")

        elif slide_layout == prs.slide_layouts[1]:
            slide.placeholders[10].text = slide_data.get("title", "")
            slide.placeholders[11].text = slide_data.get("undertitle", "")
            # random_image_path = get_random_image('/content/images')
            slide.placeholders[12].insert_picture(image_stream)

        elif slide_layout == prs.slide_layouts[2]:
            slide.placeholders[10].text = slide_data.get("title", "")
            slide.placeholders[11].text = slide_data.get("object", "")

    user_folder = os.path.join(WORK_FOLDER, str(user_id))
    if not os.path.exists(user_folder):
        os.makedirs(user_folder)

    presentation_file_name = "your_deck.pptx"
    presentation_path = os.path.join(user_folder, presentation_file_name)
    prs.save(presentation_path)

    return presentation_path
